import pdfplumber
from docx import Document
from pptx import Presentation
import xmindparser
from pdfminer.high_level import extract_text
import pytesseract
from PIL import Image
import argparse
import sys
import os


def parse_pdf(filepath):
    """解析PDF文件,优先使用pdfminer,失败则尝试pdfplumber或OCR"""
    text = extract_text(filepath)
    if text:
        return text

    # 尝试使用pdfplumber逐页提取
    text = []
    try:
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text)
                else:
                    # OCR 提取图片文字
                    image = page.to_image().original
                    ocr_text = pytesseract.image_to_string(image)
                    text.append(ocr_text)
        return '\n'.join(text)
    except Exception as e:
        print(f"PDF解析错误: {str(e)}", file=sys.stderr)
        return ''


def parse_docx(filepath):
    """解析Word文档"""
    try:
        doc = Document(filepath)
        return '\n'.join([para.text for para in doc.paragraphs])
    except Exception as e:
        print(f"DOCX解析错误: {str(e)}", file=sys.stderr)
        return ''


def parse_pptx(filepath):
    """解析PPT文档""" 
    try:
        prs = Presentation(filepath)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    except Exception as e:
        print(f"PPTX解析错误: {str(e)}", file=sys.stderr)
        return ''


def parse_xmind(filepath):
    """解析XMind文件"""
    try:
        content = xmindparser.xmind_to_dict(filepath)
        return str(content)
    except Exception as e:
        print(f"XMind解析错误: {str(e)}", file=sys.stderr)
        return ''


def parse_txt(filepath):
    """解析TXT文件"""
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        print(f"TXT解析错误: {str(e)}", file=sys.stderr)
        return ''


def parse_md(filepath):
    """解析Markdown文件"""
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        print(f"MD解析错误: {str(e)}", file=sys.stderr)
        return ''


def parse_file(filepath):
    """根据文件类型自动选择解析器"""
    ext = os.path.splitext(filepath)[1].lower()
    if ext == '.pdf':
        return parse_pdf(filepath)
    elif ext in ['.doc', '.docx']:
        return parse_docx(filepath)
    elif ext in ['.ppt', '.pptx']:
        return parse_pptx(filepath)
    elif ext == '.xmind':
        return parse_xmind(filepath)
    elif ext in ['.txt', '.md']:
        return parse_txt(filepath)
    else:
        raise ValueError(f"不支持的格式: {ext}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description='文件解析工具')
    parser.add_argument('filepath', help='需要解析的文件路径')
    args = parser.parse_args()
    try:
        print(parse_file(args.filepath))
    except Exception as e:
        print(f"错误: {str(e)}", file=sys.stderr)
        sys.exit(1)
