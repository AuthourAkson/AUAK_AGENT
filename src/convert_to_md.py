import os
import pdfplumber
import pandas as pd
from pptx import Presentation
from docx import Document
from xmindparser import xmind_to_dict

PRE_KNOWLEDGE_DIR = os.path.join(os.path.dirname(__file__), '..', 'pre-knowledge')
KNOWLEDGE_DIR = os.path.join(os.path.dirname(__file__), '..', 'knowledge')

def convert_docx_to_md(path):
    doc = Document(path)
    return '\n\n'.join(p.text for p in doc.paragraphs if p.text.strip())

def convert_pdf_to_md(path):
    text = ''
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() + '\n'
    return text.strip()

def convert_txt_to_md(path):
    with open(path, 'r', encoding='utf-8') as f:
        return f.read()

def convert_csv_to_md(path):
    df = pd.read_csv(path)
    return df.to_markdown(index=False)

def convert_ppt_to_md(path):
    prs = Presentation(path)
    content = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                slide_text.append(shape.text)
        content.append('\n'.join(slide_text))
    return '\n\n---\n\n'.join(content)

def convert_xmind_to_md(path):
    xmind_data = xmind_to_dict(path)
    md_lines = []

    def walk_topic(topic, level=1):
        prefix = "#" * level
        md_lines.append(f"{prefix} {topic.get('title', '')}")
        for child in topic.get('topics', []):
            walk_topic(child, level + 1)

    for sheet in xmind_data:
        md_lines.append(f"# {sheet.get('title', '')}")
        if 'topic' in sheet:
            walk_topic(sheet['topic'], level=2)

    return '\n\n'.join(md_lines)

def convert_file_to_md(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    if ext == '.docx':
        return convert_docx_to_md(filepath)
    elif ext == '.pdf':
        return convert_pdf_to_md(filepath)
    elif ext == '.txt':
        return convert_txt_to_md(filepath)
    elif ext == '.csv':
        return convert_csv_to_md(filepath)
    elif ext in ['.ppt', '.pptx']:
        return convert_ppt_to_md(filepath)
    elif ext == '.xmind':
        return convert_xmind_to_md(filepath)
    else:
        return None

def main():
    for filename in os.listdir(PRE_KNOWLEDGE_DIR):
        src_path = os.path.join(PRE_KNOWLEDGE_DIR, filename)
        if os.path.isfile(src_path):
            print(f'Processing: {filename}')
            md_content = convert_file_to_md(src_path)
            if md_content:
                md_filename = os.path.splitext(filename)[0] + '.md'
                dest_path = os.path.join(KNOWLEDGE_DIR, md_filename)
                with open(dest_path, 'w', encoding='utf-8') as f:
                    f.write(md_content)
                print(f'✓ Converted and saved to: {md_filename}')
            else:
                print(f'✗ Unsupported file type: {filename}')

if __name__ == '__main__':
    main()
